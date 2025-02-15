import streamlit as st
import magic
import filetype
import PyPDF2
import docx
import pandas as pd
from pptx import Presentation
from collections import Counter
from py3langid.langid import LanguageIdentifier, MODEL_FILE
import os

import magic
import filetype
import PyPDF2
import docx
import pandas as pd
from pptx import Presentation
from collections import Counter
from py3langid.langid import LanguageIdentifier, MODEL_FILE

def detect_file_type(file_path):
    kind = filetype.guess(file_path)
    if kind:
        return kind.extension
    mime = magic.Magic(mime=True)
    return mime.from_file(file_path)

def extract_text(file_path, file_type):
    text = ""
    try:
        if file_type in ["pdf"]:
            with open(file_path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    text += page.extract_text() or ""
        elif file_type in ["docx"]:
            doc = docx.Document(file_path)
            text = "\n".join([para.text for para in doc.paragraphs])
        elif file_type in ["xlsx", "csv"]:
            df = pd.read_excel(file_path) if file_type == "xlsx" else pd.read_csv(file_path)
            text = df.to_string()
        elif file_type in ["pptx"]:
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        elif file_type in ["txt"]:
            with open(file_path, "r", encoding="utf-8") as f:
                text = f.read()
        else:
            with open(file_path, "rb") as f:
                text = f.read().decode(errors='ignore')
    except Exception as e:
        return f"Error extracting text: {str(e)}"
    return text

def detect_language_with_langid(line):
    # Initialize the language identifier with the model
    identifier = LanguageIdentifier.from_pickled_model(MODEL_FILE, norm_probs=True)
    lang, prob = identifier.classify(line)
    return lang, prob

def clean_text(text):
   
    return text.strip()

def detect_languages(text):
    lines = text.split('\n')
    detected_languages = set()  # Use set to avoid duplicates
    for line in lines:
        line = clean_text(line)
        if line:  # Skip empty lines
            lang, prob = detect_language_with_langid(line)
            # You can apply a threshold on probability to filter out less confident predictions
            if prob > 0.5:  # Example: filter out languages with less than 50% confidence
                detected_languages.add(lang)
    return list(detected_languages)  # Return list of detected languages

def main(file_path):
    file_type = detect_file_type(file_path)
    text = extract_text(file_path, file_type)
    languages = detect_languages(text)
    return {"File Type": file_type, "Languages": languages}




def main():
    st.title("File Type and Language Detector")

    uploaded_file = st.file_uploader("Choose a file", type=["pdf", "docx", "xlsx", "csv", "pptx", "txt"])

    if uploaded_file is not None:
        # Save the uploaded file temporarily
        with open("temp_file", "wb") as f:
            f.write(uploaded_file.getbuffer())

        try:
            file_type = detect_file_type("temp_file")
            text = extract_text("temp_file", file_type)
            languages = detect_languages(text)

            st.success("File processed successfully!")
            st.write("File Type:", file_type)
            st.write("Detected Languages:", ", ".join(languages))

            # Display a sample of the extracted text
            st.subheader("Sample of Extracted Text")
            st.text(text[:500] + "..." if len(text) > 500 else text)

        except Exception as e:
            st.error(f"An error occurred: {str(e)}")

        finally:
            # Remove the temporary file
            os.remove("temp_file")

if __name__ == "__main__":
    main()
